import os
import re
import json
import argparse
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from openai import OpenAI

def delete_slide(presentation, index):
    """Deletes a slide from the presentation by its index."""
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])

def parse_markdown(md_content):
    """Extracts text and identifies image paths found in the report."""
    images = re.findall(r'!\[.*?\]\((.*?)\)', md_content)
    return md_content, images

def estimate_text_capacity(placeholder):
    """Estimates how many bullet points and characters a placeholder can hold."""
    width_inches = placeholder.width / 914400
    height_inches = placeholder.height / 914400
    
    # More aggressive estimation - typical slides can hold more text
    # Assume ~10pt font, ~8-10 lines per inch vertically
    # Each bullet takes about 0.3-0.35 inches with spacing
    max_bullets = max(3, int(height_inches / 0.3))
    
    # Estimate max chars per bullet based on width
    # At 10-12pt, you can fit roughly 15-18 chars per inch
    max_chars_per_bullet = max(50, int(width_inches * 15))
    
    # Cap at reasonable maximums to maintain readability
    max_bullets = min(max_bullets, 10)
    max_chars_per_bullet = min(max_chars_per_bullet, 120)
    
    return max_bullets, max_chars_per_bullet

def get_template_layouts(pptx_path):
    """Extracts detailed metadata from the template to help the LLM select layouts."""
    prs = Presentation(pptx_path)
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        text_capacity = None
        
        for p in layout.placeholders:
            ph_info = {
                "name": p.name,
                "type": str(p.placeholder_format.type),
                "left": p.left,
                "top": p.top,
                "width": p.width,
                "height": p.height
            }
            placeholders.append(ph_info)
            
            if p.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                max_bullets, max_chars = estimate_text_capacity(p)
                text_capacity = {
                    "max_bullets": max_bullets,
                    "max_chars_per_bullet": max_chars
                }
        
        layouts.append({
            "index": i,
            "name": layout.name,
            "placeholders": placeholders,
            "placeholder_count": len(placeholders),
            "text_capacity": text_capacity
        })
    return layouts

def auto_fit_text(text_frame, max_font_size=18, min_font_size=10):
    """Automatically adjusts font size to fit content in text frame."""
    if not text_frame.text.strip():
        return
    
    for size in range(max_font_size, min_font_size - 1, -1):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size)
        
        if text_frame.text:
            break

def summarize_and_map_with_llm(md_content, layouts, image_list, api_key):
    """Consults the LLM to structure the report and pick the best template layouts."""
    client = OpenAI(api_key=api_key)
    
    layout_descriptions = []
    for l in layouts:
        ph_types = [p['type'] for p in l['placeholders']]
        has_picture = any('PICTURE' in t for t in ph_types)
        has_body = any('BODY' in t or 'OBJECT' in t for t in ph_types)
        
        desc = f"- Index {l['index']}: '{l['name']}'"
        
        if l.get('text_capacity'):
            cap = l['text_capacity']
            desc += f" [Max: {cap['max_bullets']} bullets, {cap['max_chars_per_bullet']} chars/bullet]"
        
        if has_picture:
            desc += " [IMAGES]"
        if has_body:
            desc += " [TEXT]"
            
        layout_descriptions.append(desc)
    
    layout_context = "\n".join(layout_descriptions)
    
    prompt = f"""
    You are a professional presentation designer. Convert the following report into a slide deck.
    
    AVAILABLE TEMPLATE LAYOUTS (with text capacity limits):
    {layout_context}
    
    AVAILABLE IMAGES:
    {", ".join(image_list) if image_list else "No images available"}
    
    CRITICAL INSTRUCTIONS:
    1. **RESPECT LAYOUT CAPACITY**: Each layout shows its max bullets and chars per bullet
    2. **USE THE FULL CAPACITY** - Don't be overly conservative, utilize the available space
    3. If content naturally fits in fewer bullets, that's fine, but don't artificially split when you have room
    4. Create a title slide first
    5. For each slide, choose a layout with adequate capacity for your content
    6. Only split into multiple slides if content truly exceeds the layout's stated capacity
    7. For images, prefer layouts marked [IMAGES]
    8. Aim for informative, substantive content that uses available space effectively
    
    RETURN FORMAT (JSON only):
    {{
      "slides": [
        {{
          "title": "slide title (under 45 chars)",
          "bullets": ["bullet 1", "bullet 2", ...],
          "image_path": "path/to/image.png or null",
          "layout_index": 0,
          "notes": "optional context"
        }}
      ]
    }}
    
    REPORT:
    {md_content}
    """

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": "You are a PowerPoint expert. Output ONLY valid JSON. Respect layout capacity limits but use the available space effectively. Create comprehensive, informative slides."},
            {"role": "user", "content": prompt}
        ],
        response_format={"type": "json_object"}
    )
    return json.loads(response.choices[0].message.content).get('slides', [])

def fit_image_to_placeholder(placeholder, image_path):
    """Inserts image while maintaining aspect ratio within placeholder bounds."""
    from PIL import Image
    
    ph_width = placeholder.width
    ph_height = placeholder.height
    ph_left = placeholder.left
    ph_top = placeholder.top
    
    with Image.open(image_path) as img:
        img_width, img_height = img.size
        img_ratio = img_width / img_height
        ph_ratio = ph_width / ph_height
        
        if img_ratio > ph_ratio:
            new_width = ph_width
            new_height = int(ph_width / img_ratio)
        else:
            new_height = ph_height
            new_width = int(ph_height * img_ratio)
        
        left = ph_left + (ph_width - new_width) // 2
        top = ph_top + (ph_height - new_height) // 2
        
        return left, top, new_width, new_height

def split_slide_if_needed(slide_info, layouts):
    """Splits a slide into multiple slides based on the chosen layout's actual capacity."""
    bullets = slide_info.get('bullets', [])
    layout_idx = slide_info.get('layout_index', 1)
    
    if layout_idx < len(layouts) and layouts[layout_idx].get('text_capacity'):
        capacity = layouts[layout_idx]['text_capacity']
        max_bullets = capacity['max_bullets']
        max_chars = capacity['max_chars_per_bullet']
    else:
        max_bullets = 6
        max_chars = 80
    
    # Only split if we significantly exceed capacity (add 20% buffer)
    needs_split = len(bullets) > int(max_bullets * 1.2)
    
    # Check for excessively long bullets (30% over limit)
    for bullet in bullets:
        if len(str(bullet)) > int(max_chars * 1.3):
            needs_split = True
            break
    
    if not needs_split:
        return [slide_info]
    
    result_slides = []
    chunks = []
    current_chunk = []
    
    for bullet in bullets:
        bullet_text = str(bullet)
        
        if len(bullet_text) > int(max_chars * 1.3):
            if '. ' in bullet_text:
                sentences = bullet_text.split('. ')
                for sent in sentences:
                    if not sent.endswith('.'):
                        sent += '.'
                    if len(current_chunk) >= max_bullets:
                        chunks.append(current_chunk)
                        current_chunk = [sent.strip()]
                    else:
                        current_chunk.append(sent.strip())
            else:
                if current_chunk:
                    chunks.append(current_chunk)
                    current_chunk = []
                chunks.append([bullet_text])
        else:
            if len(current_chunk) >= max_bullets:
                chunks.append(current_chunk)
                current_chunk = [bullet_text]
            else:
                current_chunk.append(bullet_text)
    
    if current_chunk:
        chunks.append(current_chunk)
    
    base_title = slide_info.get('title', 'Slide')
    total_parts = len(chunks)
    
    for i, chunk in enumerate(chunks):
        new_slide = slide_info.copy()
        if total_parts > 1:
            clean_title = re.sub(r'\s*\(\d+/\d+\)\s*$', '', base_title)
            clean_title = re.sub(r'\s*-\s*Part\s+\d+\s*$', '', clean_title, flags=re.IGNORECASE)
            new_slide['title'] = f"{clean_title} ({i+1}/{total_parts})"
        new_slide['bullets'] = chunk
        if i > 0:
            new_slide['image_path'] = None
        result_slides.append(new_slide)
    
    return result_slides

def create_presentation(slides_data, template_path, output_path, md_dir, layouts):
    """Constructs the presentation with improved text fitting and image placement."""
    prs = Presentation(template_path)
    
    num_original_slides = len(prs.slides)
    print(f"Template contains {num_original_slides} sample slides. They will be removed after generation.")

    processed_slides = []
    for slide_info in slides_data:
        split_slides = split_slide_if_needed(slide_info, layouts)
        processed_slides.extend(split_slides)
    
    if len(processed_slides) > len(slides_data):
        print(f"Note: Split {len(slides_data)} slides into {len(processed_slides)} slides for better readability")

    for i, slide_info in enumerate(processed_slides):
        l_idx = slide_info.get('layout_index', 1)
        if l_idx >= len(prs.slide_layouts):
            l_idx = 1
            
        layout = prs.slide_layouts[l_idx]
        slide = prs.slides.add_slide(layout)
        print(f"Creating Slide {i+1}: {slide_info.get('title', 'Untitled')}")

        if slide.shapes.title:
            slide.shapes.title.text = slide_info.get('title', 'Slide')
            if slide.shapes.title.has_text_frame:
                auto_fit_text(slide.shapes.title.text_frame, max_font_size=32, min_font_size=20)

        text_placeholder = next(
            (s for s in slide.placeholders if s.placeholder_format.type in 
             [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.OBJECT]), 
            None
        )
        
        if text_placeholder and slide_info.get('bullets'):
            tf = text_placeholder.text_frame
            tf.clear()
            tf.word_wrap = True
            
            for bullet_text in slide_info.get('bullets', []):
                p = tf.add_paragraph()
                p.text = str(bullet_text)
                p.level = 0
            
            auto_fit_text(tf, max_font_size=18, min_font_size=10)

        img_rel = slide_info.get('image_path')
        if img_rel:
            full_img_path = os.path.normpath(os.path.join(md_dir, img_rel))
            if os.path.exists(full_img_path):
                try:
                    pic_placeholder = next(
                        (s for s in slide.placeholders if s.placeholder_format.type == PP_PLACEHOLDER.PICTURE), 
                        None
                    )
                    
                    if pic_placeholder:
                        left, top, width, height = fit_image_to_placeholder(pic_placeholder, full_img_path)
                        sp = pic_placeholder._element
                        sp.getparent().remove(sp)
                        slide.shapes.add_picture(full_img_path, left, top, width=width, height=height)
                    else:
                        slide.shapes.add_picture(
                            full_img_path, 
                            Inches(6), 
                            Inches(1.5), 
                            width=Inches(3.5)
                        )
                except Exception as e:
                    print(f"Warning: Could not insert image {img_rel}: {e}")

    print(f"Cleaning up {num_original_slides} sample slides...")
    for _ in range(num_original_slides):
        delete_slide(prs, 0)

    try:
        prs.save(output_path)
        print(f"\n✓ SUCCESS: Created presentation at {output_path}")
    except PermissionError:
        print(f"\n❌ ERROR: Cannot save to '{output_path}'")
        print("The file may have been opened during processing.")
        print("Please close it and run the script again.")
        sys.exit(1)

def validate_files(md_path, pptx_path):
    if not os.path.exists(md_path):
        print(f"Error: Markdown file not found: {md_path}")
        sys.exit(1)
    if not os.path.exists(pptx_path):
        print(f"Error: PowerPoint template not found: {pptx_path}")
        sys.exit(1)

def check_output_writable(output_path):
    """Check if output file can be written (not open in another program)."""
    if os.path.exists(output_path):
        try:
            with open(output_path, 'a'):
                pass
            return True
        except PermissionError:
            return False
    return True

def main():
    parser = argparse.ArgumentParser(description="Convert Markdown to PowerPoint with AI-powered layout selection")
    parser.add_argument("--md", required=True, help="Path to input Markdown file")
    parser.add_argument("--pptx", required=True, help="Path to PowerPoint template")
    parser.add_argument("--api_key", help="OpenAI API key (or set OPENAI_API_KEY env var)")
    parser.add_argument("--output", help="Output file path (default: adds '_output' to template name)")
    args = parser.parse_args()

    validate_files(args.md, args.pptx)
    
    api_key = args.api_key or os.getenv("OPENAI_API_KEY")
    if not api_key:
        print("Error: OpenAI API key required. Use --api_key or set OPENAI_API_KEY environment variable.")
        sys.exit(1)
    
    md_abs = os.path.abspath(args.md)
    md_dir = os.path.dirname(md_abs)
    
    if args.output:
        output_path = args.output
    else:
        base, ext = os.path.splitext(args.pptx)
        output_path = f"{base}_output{ext}"
    
    if not check_output_writable(output_path):
        print(f"\n❌ ERROR: Cannot write to '{output_path}'")
        print("The file may be open in PowerPoint or another program.")
        print("Please close the file and try again.\n")
        sys.exit(1)

    with open(md_abs, 'r', encoding='utf-8') as f:
        content = f.read()

    _, images = parse_markdown(content)
    layouts = get_template_layouts(args.pptx)
    
    print("Analyzing content and generating slide structure...")
    slides_json = summarize_and_map_with_llm(content, layouts, images, api_key)
    
    print(f"Creating presentation with {len(slides_json)} slides...")
    create_presentation(slides_json, args.pptx, output_path, md_dir, layouts)

if __name__ == "__main__":
    main()